#!/usr/bin/env python3
"""
Document Extractor for Vibe Coding Workflow

PPT, PDF, MD 파일에서 텍스트와 이미지를 추출하여
AI 문서 변환 워크플로우에 사용할 수 있는 형태로 출력한다.

사용법:
    python tools/extract-doc.py <파일경로> --service <서비스명> --version <버전>

예시:
    python tools/extract-doc.py ~/기획서.pptx --service auth --version v1.0.0
    python tools/extract-doc.py ~/설계서.pdf --service payment --version v0.1.0
"""

import argparse
import os
import re
import shutil
import sys
from pathlib import Path


def extract_pptx(file_path: str, assets_dir: str) -> str:
    from pptx import Presentation
    from pptx.enum.shapes import MSO_SHAPE_TYPE

    prs = Presentation(file_path)
    sections = []
    image_count = 0

    for slide_idx, slide in enumerate(prs.slides, 1):
        slide_lines = [f"## 슬라이드 {slide_idx}"]

        if slide.has_notes_slide and slide.notes_slide.notes_text_frame:
            notes = slide.notes_slide.notes_text_frame.text.strip()
            if notes:
                slide_lines.append(f"\n> 발표자 노트: {notes}")

        texts = []
        for shape in slide.shapes:
            if shape.has_text_frame:
                for paragraph in shape.text_frame.paragraphs:
                    text = paragraph.text.strip()
                    if text:
                        texts.append(text)

            if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
                image_count += 1
                image = shape.image
                ext = image.content_type.split("/")[-1]
                if ext == "jpeg":
                    ext = "jpg"
                filename = f"slide-{slide_idx:02d}-img-{image_count:03d}.{ext}"
                image_path = os.path.join(assets_dir, filename)

                with open(image_path, "wb") as f:
                    f.write(image.blob)

                texts.append(f"\n![슬라이드 {slide_idx} 이미지](assets/{filename})")

            if shape.has_table:
                table = shape.table
                header = [cell.text.strip() for cell in table.rows[0].cells]
                texts.append("")
                texts.append("| " + " | ".join(header) + " |")
                texts.append("| " + " | ".join("---" for _ in header) + " |")
                for row in table.rows[1:]:
                    cells = [cell.text.strip() for cell in row.cells]
                    texts.append("| " + " | ".join(cells) + " |")

        if texts:
            slide_lines.append("")
            slide_lines.extend(texts)

        sections.append("\n".join(slide_lines))

    return "\n\n---\n\n".join(sections)


def extract_pdf(file_path: str, assets_dir: str) -> str:
    import fitz

    doc = fitz.open(file_path)
    sections = []
    image_count = 0

    for page_idx in range(len(doc)):
        page = doc[page_idx]
        page_num = page_idx + 1
        page_lines = [f"## 페이지 {page_num}"]

        text = page.get_text().strip()
        if text:
            page_lines.append("")
            page_lines.append(text)

        images = page.get_images(full=True)
        for img in images:
            image_count += 1
            xref = img[0]
            try:
                base_image = doc.extract_image(xref)
                image_bytes = base_image["image"]
                ext = base_image["ext"]

                if len(image_bytes) < 5000:
                    continue

                filename = f"page-{page_num:02d}-img-{image_count:03d}.{ext}"
                image_path = os.path.join(assets_dir, filename)

                with open(image_path, "wb") as f:
                    f.write(image_bytes)

                page_lines.append(f"\n![페이지 {page_num} 이미지](assets/{filename})")
            except Exception:
                page_lines.append(f"\n[이미지 추출 실패 - 페이지 {page_num}, xref {xref}]")

        sections.append("\n".join(page_lines))

    doc.close()
    return "\n\n---\n\n".join(sections)


def extract_md(file_path: str, assets_dir: str) -> str:
    with open(file_path, "r", encoding="utf-8") as f:
        content = f.read()

    source_dir = os.path.dirname(os.path.abspath(file_path))
    image_count = 0
    image_pattern = r"!\[([^\]]*)\]\(([^)]+)\)"

    def copy_and_rewrite(match):
        nonlocal image_count
        alt_text = match.group(1)
        img_path = match.group(2)

        if img_path.startswith("http"):
            return match.group(0)

        source_image = os.path.join(source_dir, img_path)
        if os.path.exists(source_image):
            image_count += 1
            ext = os.path.splitext(img_path)[1]
            filename = f"md-img-{image_count:03d}{ext}"
            dest_path = os.path.join(assets_dir, filename)
            shutil.copy2(source_image, dest_path)
            return f"![{alt_text}](assets/{filename})"

        return match.group(0)

    content = re.sub(image_pattern, copy_and_rewrite, content)
    return content


EXTRACTORS = {
    ".pptx": extract_pptx,
    ".pdf": extract_pdf,
    ".md": extract_md,
}


def generate_output(source_file: str, extracted_text: str, assets_dir: str) -> str:
    asset_files = sorted(os.listdir(assets_dir)) if os.path.exists(assets_dir) else []
    image_files = [f for f in asset_files if f.lower().endswith((".png", ".jpg", ".jpeg", ".gif", ".svg", ".webp"))]

    lines = [
        "# 문서 추출 결과",
        "",
        f"- **원본 파일**: `{os.path.basename(source_file)}`",
        f"- **추출된 이미지**: {len(image_files)}개 → `{assets_dir}`",
        "",
        "---",
        "",
        "# 추출된 내용",
        "",
        extracted_text,
    ]

    if image_files:
        lines.extend([
            "",
            "---",
            "",
            "# 추출된 이미지 목록",
            "",
            "| # | 파일명 | 용도 (수동 분류 필요) |",
            "|---|--------|---------------------|",
        ])
        for i, img in enumerate(image_files, 1):
            lines.append(f"| {i} | `{img}` | |")

        lines.extend([
            "",
            "> 위 이미지들을 확인하고, `02-screen-spec.md` 작성 시 각 화면(S-ID)에 매핑하세요.",
        ])

    return "\n".join(lines)


def main():
    parser = argparse.ArgumentParser(
        description="PPT/PDF/MD에서 텍스트와 이미지를 추출합니다.",
        epilog="추출 결과를 Cursor 채팅에 붙여넣고 @docs/_templates/ 와 함께 문서 변환을 요청하세요.",
    )
    parser.add_argument("file", help="추출할 파일 경로 (.pptx, .pdf, .md)")
    parser.add_argument("--service", required=True, help="서비스명 (예: auth, payment)")
    parser.add_argument("--version", required=True, help="버전 (예: v1.0.0)")
    parser.add_argument("--output", help="결과 저장 경로 (미지정 시 stdout 출력)")

    args = parser.parse_args()

    file_path = os.path.abspath(args.file)
    if not os.path.exists(file_path):
        print(f"오류: 파일을 찾을 수 없습니다 - {file_path}", file=sys.stderr)
        sys.exit(1)

    ext = os.path.splitext(file_path)[1].lower()
    if ext not in EXTRACTORS:
        print(f"오류: 지원하지 않는 형식입니다 - {ext}", file=sys.stderr)
        print(f"지원 형식: {', '.join(EXTRACTORS.keys())}", file=sys.stderr)
        sys.exit(1)

    project_root = Path(__file__).resolve().parent.parent
    assets_dir = project_root / "docs" / args.service / args.version / "assets"
    assets_dir.mkdir(parents=True, exist_ok=True)

    print(f"추출 중: {os.path.basename(file_path)} → docs/{args.service}/{args.version}/assets/", file=sys.stderr)

    extracted = EXTRACTORS[ext](file_path, str(assets_dir))
    output = generate_output(file_path, extracted, str(assets_dir))

    if args.output:
        output_path = os.path.abspath(args.output)
        with open(output_path, "w", encoding="utf-8") as f:
            f.write(output)
        print(f"완료: {output_path}", file=sys.stderr)
    else:
        print(output)

    asset_count = len(list(assets_dir.glob("*"))) if assets_dir.exists() else 0
    print(f"\n추출 완료: 이미지 {asset_count}개 저장됨 → {assets_dir}", file=sys.stderr)
    print("다음 단계: 추출 결과를 Cursor 채팅에 전달하여 문서 변환을 요청하세요.", file=sys.stderr)


if __name__ == "__main__":
    main()
